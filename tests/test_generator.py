"""
Tests for NeetSlides PPTX Generator.
"""

from pathlib import Path
import tempfile

import pytest

from neetslides.generator import create_presentation, generate_pptx, add_slide_from_data
from neetslides.models import BoundingBox, DocumentData, SlideData, TextBlock


def create_test_analyzed_document() -> DocumentData:
    """Create a test document with semantic analysis already applied."""
    slides = []
    
    for page_num in range(2):
        blocks = [
            TextBlock(
                text=f"Slide {page_num + 1} Title",
                font_size=24.0,
                font_name="Arial-Bold",
                bbox=BoundingBox(x0=50, y0=50, x1=500, y1=80),
                page_num=page_num,
                semantic_type="title",
            ),
            TextBlock(
                text="First bullet point",
                font_size=12.0,
                font_name="Arial",
                bbox=BoundingBox(x0=50, y0=120, x1=500, y1=140),
                page_num=page_num,
                semantic_type="body",
                indent_level=0,
            ),
            TextBlock(
                text="Sub-bullet point",
                font_size=12.0,
                font_name="Arial",
                bbox=BoundingBox(x0=80, y0=150, x1=500, y1=170),
                page_num=page_num,
                semantic_type="body",
                indent_level=1,
            ),
        ]
        
        slide = SlideData(
            page_num=page_num,
            width=612,
            height=792,
            text_blocks=blocks,
            title=f"Slide {page_num + 1} Title",
            body_blocks=[b for b in blocks if b.semantic_type == "body"],
        )
        slides.append(slide)
    
    return DocumentData(
        source_path=Path("test.pdf"),
        total_pages=2,
        slides=slides,
    )


class TestPresentationCreation:
    """Tests for presentation creation."""
    
    def test_create_presentation(self):
        """Test that a presentation can be created."""
        prs = create_presentation()
        assert prs is not None
        assert len(prs.slides) == 0


class TestSlideGeneration:
    """Tests for slide generation."""
    
    def test_add_slide(self):
        """Test adding a slide with content."""
        prs = create_presentation()
        doc = create_test_analyzed_document()
        
        add_slide_from_data(prs, doc.slides[0])
        
        assert len(prs.slides) == 1


class TestPPTXGeneration:
    """Tests for full PPTX generation."""
    
    def test_generate_pptx(self):
        """Test generating a complete PPTX file."""
        doc = create_test_analyzed_document()
        
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            result = generate_pptx(doc, output_path)
            
            assert result.exists()
            assert result.suffix == ".pptx"
            
            # Verify file is valid by opening it
            from pptx import Presentation
            prs = Presentation(result)
            assert len(prs.slides) == 2
