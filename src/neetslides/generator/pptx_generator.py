"""
NeetSlides PPTX Generator - Create editable PowerPoint files.

Uses python-pptx to generate PPTX files with proper semantic
placeholders for titles and body content.
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from neetslides.models import DocumentData, SlideData, TextBlock


# Standard slide layouts in python-pptx
LAYOUT_TITLE_SLIDE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_SECTION_HEADER = 2
LAYOUT_TWO_CONTENT = 3
LAYOUT_COMPARISON = 4
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6


def create_presentation() -> Presentation:
    """Create a new blank presentation with standard dimensions."""
    prs = Presentation()
    # Standard 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide_from_data(prs: Presentation, slide_data: SlideData) -> None:
    """
    Add a slide to the presentation from analyzed slide data.
    
    Uses proper placeholders for titles and content to ensure
    theme reflow capability.
    """
    # Choose layout based on content
    if slide_data.title and slide_data.body_blocks:
        layout = prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT]
    elif slide_data.title:
        layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    else:
        layout = prs.slide_layouts[LAYOUT_BLANK]
    
    slide = prs.slides.add_slide(layout)
    
    # Set title if present
    if slide_data.title and slide.shapes.title:
        slide.shapes.title.text = slide_data.title
    
    # Add body content
    if slide_data.body_blocks:
        _add_body_content(slide, layout, slide_data.body_blocks)


def _add_body_content(slide, layout, body_blocks: list[TextBlock]) -> None:
    """Add body content to slide, using placeholder if available."""
    # Try to find the content placeholder
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            content_placeholder = shape
            break
    
    if content_placeholder is not None:
        # Use the placeholder's text frame
        tf = content_placeholder.text_frame
        tf.clear()  # Clear any default text
        
        for i, block in enumerate(body_blocks):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = block.text
            p.level = min(block.indent_level, 8)  # Max 9 levels (0-8)
            
            # Apply font settings
            for run in p.runs:
                run.font.size = Pt(18)  # Standard body size
    else:
        # Fallback: add text box
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(12)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, block in enumerate(body_blocks):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = block.text
            
            # Simulate indent with left margin
            if block.indent_level > 0:
                p.level = min(block.indent_level, 8)


def generate_pptx(
    doc: DocumentData,
    output_path: Path,
    normalize_fonts: bool = True
) -> Path:
    """
    Generate a PPTX file from analyzed document data.
    
    Args:
        doc: DocumentData with semantic analysis applied
        output_path: Where to save the PPTX file
        normalize_fonts: Whether to normalize fonts for readability
        
    Returns:
        Path to the generated PPTX file
    """
    prs = create_presentation()
    
    for slide_data in doc.slides:
        add_slide_from_data(prs, slide_data)
    
    # Ensure output directory exists
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    prs.save(output_path)
    return output_path


def convert_pdf_to_pptx(
    pdf_path: Path,
    output_path: Optional[Path] = None,
    verbose: bool = False
) -> Path:
    """
    High-level function to convert a PDF to PPTX.
    
    Args:
        pdf_path: Path to input PDF
        output_path: Path for output PPTX (defaults to same name as PDF)
        verbose: Print progress information
        
    Returns:
        Path to the generated PPTX file
    """
    from neetslides.parser import parse_pdf
    from neetslides.heuristics import analyze_document
    
    pdf_path = Path(pdf_path)
    
    if output_path is None:
        output_path = pdf_path.with_suffix(".pptx")
    
    if verbose:
        print(f"Parsing PDF: {pdf_path}")
    
    # Parse PDF
    doc = parse_pdf(pdf_path)
    
    if verbose:
        print(f"  Found {doc.total_pages} pages")
    
    # Analyze semantics
    if verbose:
        print("Analyzing slide structure...")
    
    analyzed_doc = analyze_document(doc)
    
    if verbose:
        for slide in analyzed_doc.slides:
            print(f"  Page {slide.page_num + 1}: {slide.title or '(no title)'}")
    
    # Generate PPTX
    if verbose:
        print(f"Generating PPTX: {output_path}")
    
    result = generate_pptx(analyzed_doc, output_path)
    
    if verbose:
        print("Done!")
    
    return result
